"""
Innovation Intelligence Platform
US Personal Care Division — Genommalab
"""

import streamlit as st
import anthropic
import os
import base64
import pandas as pd
from datetime import datetime

# ─────────────────────────────────────────────────────────────────
# BRAND KNOWLEDGE BASE
# ─────────────────────────────────────────────────────────────────

BRAND_PROFILES = {
    "Cicatricure": {
        "category": "Skin Care",
        "icon": "✨",
        "price_range": "$7–$12",
        "description": (
            "Heritage anti-aging brand. Known for scar treatment, wrinkle reduction, and skin renewal. "
            "Core SKU: Gold Cream. Also eye contour, serums, body creams. "
            "Positioned as affordable but effective — the trusted remedy Latinas bring from Mexico."
        ),
        "consumer": "Women 35–60, first-generation Mexican-American. Want anti-aging results without complex routines or high prices. Concerns: wrinkles, dark spots, uneven tone, loss of firmness.",
        "strengths": "Heritage trust, proven in Mexico, loyal repeat buyers, clinical-feel branding",
        "weaknesses": "Limited US brand awareness outside unacculturated Hispanics, small portfolio vs. competitors",
    },
    "Goicochea": {
        "category": "Body Care",
        "icon": "💪",
        "price_range": "$6–$10",
        "description": "Iconic body care brand famous for circulation cream and anti-cellulite formulas. Strong heritage across Mexico and Latin America.",
        "consumer": "Women 30–55, working-class Latinas. Want affordable body care with real functional benefits.",
        "strengths": "Strong body care heritage, functional benefit claims, loyal consumer base",
        "weaknesses": "Limited portfolio, not well known outside Latin American immigrants",
    },
    "Teatrical": {
        "category": "Skin Care",
        "icon": "🌟",
        "price_range": "$5–$9",
        "description": "Skin brightening and lightening brand with deep cultural roots in Latin America.",
        "consumer": "Latinas 25–50 who want brighter, more even-toned skin.",
        "strengths": "Cultural relevance, unique brightening positioning, loyal niche consumer",
        "weaknesses": "Sensitive category (whitening/lightening scrutiny), limited claim modernization",
    },
    "Tio Nacho": {
        "category": "Hair Care",
        "icon": "🌿",
        "price_range": "$5–$9",
        "description": (
            "Iconic Mexican hair care brand. Known for plant-based ingredients: royal jelly, garlic, aloe vera, "
            "argan oil, and ancient Mexican herbs. Core SKUs: shampoo and conditioner across multiple lines. "
            "Enormous nostalgia factor — Mexicans recognize it from childhood."
        ),
        "consumer": "Women 25–55, Mexican-origin. Dealing with hair loss, dryness, or damage. Trust Tio Nacho for its connection to Mexican herbal traditions.",
        "strengths": "Nostalgia and emotional connection, herbal/natural positioning, trusted formula heritage",
        "weaknesses": "Perceived as old brand, packaging dated vs. modern natural brands, limited recent innovation",
    },
    "Vanart": {
        "category": "Hair Care",
        "icon": "💇",
        "price_range": "$3–$6",
        "description": "Mass-market hair care brand. Wide portfolio, competes on value.",
        "consumer": "Broad Hispanic consumer, highly price-sensitive, large family households.",
        "strengths": "Price leadership, wide SKU range, broad appeal",
        "weaknesses": "No strong brand story, commodity perception, easily substituted",
    },
    "Asepxia": {
        "category": "Skin Care",
        "icon": "🧴",
        "price_range": "$5–$9",
        "description": "Acne and oily skin care specialist trusted across Latin America.",
        "consumer": "Teens and adults 15–35 with acne-prone or oily skin.",
        "strengths": "Category authority, trusted from a young age, loyal family-driven adoption",
        "weaknesses": "Dated formula perception vs. new-gen acne brands, packaging needs modernization",
    },
}

CATEGORY_BRANDS = {
    "Skin Care":  ["Cicatricure", "Teatrical", "Asepxia"],
    "Hair Care":  ["Tio Nacho", "Vanart"],
    "Body Care":  ["Goicochea"],
}

# ─────────────────────────────────────────────────────────────────
# SYSTEM PROMPT
# ─────────────────────────────────────────────────────────────────

SYSTEM_PROMPT = """You are an AI-powered Innovation Intelligence Platform for a Personal Care company operating in the United States.

You behave as a hybrid of: CPG Strategy Consultant · Data Analyst · Trend Intelligence Expert · Product Innovation Lead.

CORE OBJECTIVE: Generate clear, concise, and visually structured innovation reports that help brand teams identify:
- Growth opportunities
- Competitive gaps
- New product concepts
- Pricing strategies
- Potential suppliers

All outputs MUST be: Executive-ready · Easy to scan · Insightful but concise · Focused ONLY on the US market.

══════════════════════════════════════════════════════════════════
MANDATORY OUTPUT STRUCTURE — follow this exact order every time:
══════════════════════════════════════════════════════════════════

## 📊 1) CATEGORY SNAPSHOT — TOP 5 INSIGHTS

Five concise bullet points:
• Total category size ($) — search for current US retail figures
• 3-year growth trend (YoY %)
• YTD performance
• Key growth drivers
• Strategic insight for the selected brand

Keep it VERY concise.

## 📈 2) TOP 3 TRENDS

For each trend:
**[Trend Name]**
• What it is: 1–2 lines
• Why it matters for [selected brand]: one line
• Source: Mintel / NielsenIQ / Google Trends / TikTok / etc.

## 🏆 3) TOP 5 BRANDS

For each brand (bullet format):
**[Brand Name]** — Sales: $XM | 3Y Growth: X% | YTD: X%
• Top SKUs: [1–2 examples]
• Key strategic insight: one line

## 🥇 4) TOP 5 WINNING PRODUCTS

For each:
**[Product Name]** — [Brand] — $X.XX MSRP
• Why it's winning: one line
• Key insight: pricing / claims / channel / trend driver

## 🚀 5) INNOVATION CONCEPTS — 3 TOTAL

For EACH concept use this EXACT format:

---
### 🔹 CONCEPT [N]: [Product Name]

**🧠 INSIGHT / OPPORTUNITY**
• Market gap: one clear line
• Consumer need: one line
• Supporting trend or competitor example: one line

**🚀 CONCEPT**
• **Tagline:** "Catchy, memorable tagline"
• **Elevator pitch:** 2–3 lines max
• **Key benefits:**
  - Benefit 1
  - Benefit 2
  - Benefit 3
  - (up to 5)
• **Packaging:** Brief description aligned with brand identity

**💰 PRICE STRATEGY**
• Suggested MSRP: $X.XX
• Benchmark vs. [Competitor + Specific Product]: $X.XX
• Channel strategy: Walmart / CVS / Walgreens / Amazon

**📐 ADVANCED METRICS**
• **Opportunity size:** 💰 Small (<$10M) / 💰💰 Medium ($10M–$50M) / 💰💰💰 Large (>$50M)
• **Confidence score:** High (data-backed) / Medium (moderate signals) / Low (emerging hypothesis)
• **Speed to market:** Fast (line extension) / Medium (new formulation) / Long (new tech/claims)
• **VERDICT:** ✅ GO — brief reason / ⚠️ PROCEED WITH CAUTION — key condition / ❌ NO-GO — reason
---

## 🏭 6) SUPPLIERS (US & MEXICO) — 3 TO 5

For each:
**[Company Name]** — [City, State / Country]
• Capabilities: manufacturing / private label / ingredient supply
• Contact: [website or email]

Focus on: personal care manufacturers · private label · ingredient suppliers.

══════════════════════════════════════════════════════════════════
STYLE & CRITICAL RULES:
══════════════════════════════════════════════════════════════════
1. Bullets not paragraphs. No filler. Make it feel like a one-page strategy deck.
3. Bold all dollar figures, percentages, product names, and key verdicts.
4. Always search the web for current US market data before answering.
5. ALWAYS align concepts with the selected brand's identity, price range, and consumer profile.
6. ALWAYS consider Walmart, CVS, Walgreens, Amazon as primary channels.
7. DO NOT invent unrealistic technologies. DO NOT create vague ideas — be specific and actionable.
8. When Target Market = US Hispanic: add Hispanic-specific cultural nuances, Spanish-language behavior data.
9. When Target Market = Both: cover BOTH audiences, noting differences between them.
10. Innovation Type = Competitive Mimic + Improve: benchmark each concept against a specific winning product and explain the improvement.
11. Innovation Type = New to Market: emphasize whitespace, unmet needs, and first-mover advantage.
12. Innovation Type = Both: blend both approaches across the 3 concepts.
13. When relevant, include price tiers (good / better / best) and channel-specific insights."""

# ─────────────────────────────────────────────────────────────────
# PROMPT BUILDER
# ─────────────────────────────────────────────────────────────────

def build_prompt(category: str, brand: str, target_market: str, innovation_type: str, extra: str = "") -> str:
    p = BRAND_PROFILES.get(brand, {})

    innovation_instructions = {
        "New to Market": (
            "Focus on WHITESPACE opportunities. All 3 concepts must address unmet needs with no "
            "direct equivalent in the market. Emphasize first-mover advantage and category creation."
        ),
        "Competitive Mimic + Improve": (
            "For each of the 3 concepts, find a SPECIFIC winning competitor product, benchmark it, "
            "and show exactly how our version improves on it in formulation, claims, price, or packaging."
        ),
        "Both": (
            "Mix approaches: at least 1 concept should be true New to Market whitespace, "
            "at least 1 should be a Competitive Mimic + Improve with a clear benchmark, "
            "and 1 can blend both angles."
        ),
    }

    market_instructions = {
        "US Hispanic": (
            "Focus exclusively on the US Hispanic consumer (primarily Mexican-origin, unacculturated). "
            "Include cultural nuances, Spanish-language consumer behavior, and insights from "
            "TikTok en Español, Telemundo, beauty influencers targeting Latinas."
        ),
        "US General Market": (
            "Focus on the broader US general market consumer. Mainstream retail lens. "
            "Reference general US beauty trends, Ulta, Target, mainstream social media."
        ),
        "Both": (
            "Cover BOTH US Hispanic AND US General Market consumers. "
            "For each concept, specify which audience it primarily targets and why, "
            "and note any crossover potential."
        ),
    }

    prompt = f"""Run a full Innovation Intelligence Report with these parameters:

📂 Category: {category}
🏷️ Brand: {brand} {p.get('icon', '')} | Price Range: {p.get('price_range', 'N/A')}
🎯 Target Market: {target_market}
💡 Innovation Type: {innovation_type}

BRAND PROFILE:
• Description: {p.get('description', 'N/A')}
• Core consumer: {p.get('consumer', 'N/A')}
• Strengths: {p.get('strengths', 'N/A')}
• Weaknesses / gaps: {p.get('weaknesses', 'N/A')}

INNOVATION DIRECTIVE:
{innovation_instructions.get(innovation_type, '')}

MARKET DIRECTIVE:
{market_instructions.get(target_market, '')}

{('ADDITIONAL FOCUS: ' + extra.strip()) if extra and extra.strip() else ''}

Search the web extensively for current US market data in the {category} category before generating the report.
Then deliver the complete 6-section Innovation Intelligence Report following your mandatory output structure exactly."""

    # Inject uploaded data if available
    data_block = st.session_state.get("data_summary", "")
    if data_block:
        prompt += (
            f"\n\n{'━' * 60}\n"
            "REAL MARKET DATA (uploaded file) — prioritize these numbers over web estimates:\n\n"
            + data_block
            + f"\n{'━' * 60}"
        )

    return prompt


# ─────────────────────────────────────────────────────────────────
# FILE PARSER — IRI/Nielsen (Excel/CSV) + Documents (PDF/PPT)
# ─────────────────────────────────────────────────────────────────

def parse_data_file(uploaded_file) -> tuple:
    """Returns (df_or_None, summary_text, type_str)."""
    name = uploaded_file.name.lower()

    try:
        # ── Excel / CSV ──────────────────────────────────────────
        if name.endswith((".xlsx", ".xls", ".csv")):
            if name.endswith(".csv"):
                df = pd.read_csv(uploaded_file)
            else:
                xl = pd.ExcelFile(uploaded_file)
                if len(xl.sheet_names) == 1:
                    df = xl.parse(xl.sheet_names[0])
                else:
                    dfs = {s: xl.parse(s) for s in xl.sheet_names}
                    df = max(dfs.values(), key=lambda x: x.size)

            df.columns = [str(c).strip() for c in df.columns]
            df = df.dropna(how="all").reset_index(drop=True)

            lines = [
                "=== IRI / NIELSEN MARKET DATA ===",
                f"File: {uploaded_file.name}",
                f"Rows: {len(df):,}   Columns: {len(df.columns)}",
                f"Columns: {', '.join(df.columns.tolist())}",
                "",
            ]

            col_lower = {c.lower(): c for c in df.columns}

            def find_col(*keywords):
                for kw in keywords:
                    for cl, orig in col_lower.items():
                        if kw in cl:
                            return orig
                return None

            sales_col    = find_col("sales", "dollar", "revenue", "$")
            units_col    = find_col("unit", "volume")
            brand_col    = find_col("brand")
            prod_col     = find_col("product", "description", "item", "upc", "sku")
            chg_col      = find_col("change", "chg", "growth", "vs ya", "%")
            price_col    = find_col("price", "avg price", "$/unit")
            acv_col      = find_col("acv", "distribution", "dist")
            period_col   = find_col("period", "week", "month", "quarter", "year", "time")
            retailer_col = find_col("retailer", "channel", "account", "store", "market")

            if period_col:
                periods = df[period_col].dropna().unique()
                lines.append(f"Time periods: {', '.join(str(p) for p in periods[:8])}")
            if retailer_col:
                markets = df[retailer_col].dropna().unique()
                lines.append(f"Markets/Retailers: {', '.join(str(m) for m in markets[:10])}")
            lines.append("")

            if sales_col:
                total = pd.to_numeric(df[sales_col], errors="coerce").sum()
                lines.append(f"TOTAL CATEGORY SALES: ${total:,.0f}")
            if units_col:
                total_u = pd.to_numeric(df[units_col], errors="coerce").sum()
                lines.append(f"TOTAL CATEGORY UNITS: {total_u:,.0f}")
            lines.append("")

            if brand_col and sales_col:
                df[sales_col] = pd.to_numeric(df[sales_col], errors="coerce")
                brand_sum = (
                    df.groupby(brand_col)[sales_col]
                    .sum()
                    .sort_values(ascending=False)
                    .head(20)
                )
                total_s = brand_sum.sum()
                lines.append("BRAND SALES RANKING (Top 20):")
                for b, v in brand_sum.items():
                    share = (v / total_s * 100) if total_s > 0 else 0
                    lines.append(f"  {b}: ${v:,.0f}  ({share:.1f}% share)")
                lines.append("")

            if brand_col and chg_col:
                df[chg_col] = pd.to_numeric(df[chg_col], errors="coerce")
                chg_sum = (
                    df.groupby(brand_col)[chg_col]
                    .mean()
                    .sort_values(ascending=False)
                    .head(15)
                )
                lines.append("BRAND GROWTH — avg % change (Top 15):")
                for b, v in chg_sum.items():
                    lines.append(f"  {b}: {v:+.1f}%")
                lines.append("")

            if brand_col and price_col:
                df[price_col] = pd.to_numeric(df[price_col], errors="coerce")
                price_sum = (
                    df.groupby(brand_col)[price_col]
                    .mean()
                    .sort_values()
                    .head(15)
                )
                lines.append("AVERAGE PRICE BY BRAND (Top 15):")
                for b, v in price_sum.items():
                    lines.append(f"  {b}: ${v:.2f}")
                lines.append("")

            if brand_col and acv_col:
                df[acv_col] = pd.to_numeric(df[acv_col], errors="coerce")
                acv_sum = (
                    df.groupby(brand_col)[acv_col]
                    .mean()
                    .sort_values(ascending=False)
                    .head(15)
                )
                lines.append("ACV DISTRIBUTION BY BRAND (Top 15):")
                for b, v in acv_sum.items():
                    lines.append(f"  {b}: {v:.1f}%")
                lines.append("")

            if prod_col and sales_col:
                top_items = (
                    df.groupby(prod_col)[sales_col]
                    .sum()
                    .sort_values(ascending=False)
                    .head(15)
                )
                lines.append("TOP 15 PRODUCTS BY SALES:")
                for prod, v in top_items.items():
                    lines.append(f"  {prod}: ${v:,.0f}")
                lines.append("")

            lines.append("RAW DATA SAMPLE (first 30 rows):")
            lines.append(df.head(30).to_string(index=False))
            return df, "\n".join(lines), "table"

        # ── PDF ──────────────────────────────────────────────────
        elif name.endswith(".pdf"):
            try:
                import pypdf
                reader = pypdf.PdfReader(uploaded_file)
                parts = []
                for i, page in enumerate(reader.pages[:30]):
                    t = page.extract_text()
                    if t:
                        parts.append(f"[Page {i+1}]\n{t}")
                full_text = "\n\n".join(parts)
                return None, f"=== PDF: {uploaded_file.name} ===\n{full_text[:10000]}", "document"
            except ImportError:
                return None, f"PDF '{uploaded_file.name}' uploaded. (Install pypdf for text extraction)", "document"

        # ── PowerPoint ───────────────────────────────────────────
        elif name.endswith((".pptx", ".ppt")):
            try:
                from pptx import Presentation
                prs = Presentation(uploaded_file)
                slides = []
                for i, slide in enumerate(prs.slides):
                    lines_s = [
                        shape.text.strip()
                        for shape in slide.shapes
                        if hasattr(shape, "text") and shape.text.strip()
                    ]
                    if lines_s:
                        slides.append(f"[Slide {i+1}]\n" + "\n".join(lines_s))
                full_text = "\n\n".join(slides)
                return None, f"=== PRESENTATION: {uploaded_file.name} ===\n{full_text[:10000]}", "document"
            except ImportError:
                return None, f"PPT '{uploaded_file.name}' uploaded. (Install python-pptx for text extraction)", "document"

        else:
            return None, f"Unsupported file type: {uploaded_file.name}", "error"

    except Exception as e:
        return None, f"Could not parse file: {e}", "error"




# ─────────────────────────────────────────────────────────────────
# LOGO LOADER
# ─────────────────────────────────────────────────────────────────

def get_logo_html(height: int = 44) -> str:
    """Return an <img> tag for logo.png/svg if present, else empty string."""
    base = os.path.dirname(os.path.abspath(__file__))
    for fname in ("logo.png", "logo.svg", "logo.jpg", "logo.jpeg", "logo.webp"):
        path = os.path.join(base, fname)
        if os.path.exists(path):
            ext = fname.split(".")[-1]
            mime = "image/svg+xml" if ext == "svg" else f"image/{ext}"
            with open(path, "rb") as f:
                data = base64.b64encode(f.read()).decode()
            return (
                f'<img src="data:{mime};base64,{data}" '
                f'height="{height}" style="object-fit:contain;vertical-align:middle;'
                f'margin-right:14px" alt="Genommalab logo">'
            )
    return ""


# ─────────────────────────────────────────────────────────────────
# PAGE CONFIG
# ─────────────────────────────────────────────────────────────────

st.set_page_config(
    page_title="Innovation Intelligence Platform",
    page_icon="🔬",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# ─────────────────────────────────────────────────────────────────
# GLOBAL STYLES — DARK TECH THEME
# ─────────────────────────────────────────────────────────────────

st.markdown("""
<style>
  @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800;900&display=swap');

  /* Root & App Background */
  html, body, [data-testid="stAppViewContainer"] {
    font-family: 'Inter', sans-serif !important;
    background: #f4f7fb !important;
    color: #1e293b !important;
  }
  [data-testid="stAppViewContainer"] {
    background: #f4f7fb !important;
    min-height: 100vh;
  }
  [data-testid="stMain"] { background: transparent !important; }
  [data-testid="block-container"] { padding-top: 1.5rem !important; }

  /* Scrollbar */
  ::-webkit-scrollbar { width: 5px; }
  ::-webkit-scrollbar-track { background: #e9eef5; }
  ::-webkit-scrollbar-thumb { background: #bfcfe6; border-radius: 4px; }
  ::-webkit-scrollbar-thumb:hover { background: #0D4FA3; }

  /* Sidebar */
  [data-testid="stSidebar"] {
    background: #ffffff !important;
    border-right: 1px solid #dde5f0 !important;
    box-shadow: 2px 0 12px rgba(13,79,163,0.06) !important;
  }
  [data-testid="stSidebar"] * { color: #475569 !important; }
  [data-testid="stSidebar"] strong { color: #1e293b !important; }
  [data-testid="stSidebar"] hr { border-color: #e8edf5 !important; }
  [data-testid="stSidebar"] .stTextInput input {
    background: #f8fafc !important;
    border: 1px solid #dde5f0 !important;
    color: #1e293b !important;
    border-radius: 8px !important;
  }

  /* Animations */
  @keyframes gradientShift {
    0%   { background-position: 0% 50%; }
    50%  { background-position: 100% 50%; }
    100% { background-position: 0% 50%; }
  }
  @keyframes dot-blink {
    0%, 100% { opacity: 1; }
    50%       { opacity: 0.35; }
  }
  @keyframes button-lift {
    0%, 100% { box-shadow: 0 4px 18px rgba(13,79,163,0.28); }
    50%       { box-shadow: 0 6px 26px rgba(13,79,163,0.42); }
  }

  /* Header */
  .hub-header {
    background: linear-gradient(135deg, #0D4FA3 0%, #1a6fc4 45%, #0a3d8a 100%);
    background-size: 200% 200%;
    animation: gradientShift 7s ease infinite;
    border-radius: 16px;
    padding: 28px 36px;
    margin-bottom: 24px;
    color: white;
    border: none;
    box-shadow: 0 6px 32px rgba(13,79,163,0.22);
    position: relative;
    overflow: hidden;
  }
  .hub-header::before {
    content: "";
    position: absolute;
    top: -60px; right: -60px;
    width: 240px; height: 240px;
    background: radial-gradient(circle, rgba(255,255,255,0.08) 0%, transparent 70%);
    border-radius: 50%;
  }
  .hub-header::after {
    content: "";
    position: absolute;
    bottom: -50px; left: 25%;
    width: 180px; height: 180px;
    background: radial-gradient(circle, rgba(255,255,255,0.05) 0%, transparent 70%);
    border-radius: 50%;
  }
  .hub-header .header-badge {
    display: inline-flex;
    align-items: center;
    gap: 6px;
    background: rgba(255,255,255,0.15);
    border: 1px solid rgba(255,255,255,0.3);
    border-radius: 20px;
    padding: 4px 12px;
    font-size: 0.7rem;
    font-weight: 700;
    letter-spacing: 0.1em;
    color: #ffffff;
    margin-bottom: 12px;
  }
  .hub-header .live-dot {
    width: 6px; height: 6px;
    background: #7dd3fc;
    border-radius: 50%;
    display: inline-block;
    animation: dot-blink 1.5s ease-in-out infinite;
  }
  .hub-header h1 {
    margin: 0;
    font-size: 1.9rem;
    font-weight: 900;
    letter-spacing: -0.03em;
    color: #ffffff;
    line-height: 1.2;
  }
  .hub-header p {
    margin: 10px 0 0;
    font-size: 0.84rem;
    color: rgba(255,255,255,0.72);
    letter-spacing: 0.03em;
    font-weight: 500;
  }
  .hub-header .header-stats {
    display: flex;
    gap: 28px;
    margin-top: 18px;
  }
  .hub-header .hstat {
    font-size: 0.7rem;
    color: rgba(255,255,255,0.55);
    font-weight: 700;
    letter-spacing: 0.07em;
    text-transform: uppercase;
  }
  .hub-header .hstat span {
    display: block;
    font-size: 1.05rem;
    font-weight: 800;
    color: #ffffff;
    letter-spacing: -0.02em;
  }

  /* White cards */
  .glass-card {
    background: #ffffff;
    border: 1px solid #dde5f0;
    border-radius: 14px;
    padding: 20px 18px;
    margin-bottom: 14px;
    transition: border-color 0.2s ease, box-shadow 0.2s ease;
    box-shadow: 0 2px 10px rgba(13,79,163,0.06);
  }
  .glass-card:hover {
    border-color: #b3c9e8;
    box-shadow: 0 4px 18px rgba(13,79,163,0.10);
  }

  /* Section labels */
  .sec-label {
    font-size: 0.65rem;
    font-weight: 800;
    letter-spacing: 0.15em;
    text-transform: uppercase;
    color: #0D4FA3;
    margin: 18px 0 10px;
    display: flex;
    align-items: center;
    gap: 8px;
  }
  .sec-label::after {
    content: "";
    flex: 1;
    height: 1px;
    background: linear-gradient(90deg, #bfcfe6, transparent);
  }

  /* Selectbox */
  .stSelectbox > div > div {
    background: #ffffff !important;
    border: 1.5px solid #dde5f0 !important;
    border-radius: 10px !important;
    color: #1e293b !important;
    transition: border-color 0.2s ease !important;
  }
  .stSelectbox > div > div:hover,
  .stSelectbox > div > div:focus-within {
    border-color: #0D4FA3 !important;
    box-shadow: 0 0 0 3px rgba(13,79,163,0.1) !important;
  }
  .stSelectbox label {
    font-size: 0.75rem !important;
    font-weight: 700 !important;
    color: #64748b !important;
    letter-spacing: 0.06em !important;
    text-transform: uppercase !important;
  }
  [data-baseweb="select"] { background: transparent !important; }
  [data-baseweb="select"] * { color: #1e293b !important; }
  [data-baseweb="popover"] {
    background: #ffffff !important;
    border: 1px solid #dde5f0 !important;
    border-radius: 12px !important;
    box-shadow: 0 8px 28px rgba(13,79,163,0.12) !important;
  }
  [data-baseweb="menu"] { background: #ffffff !important; }
  [data-baseweb="option"] { background: #ffffff !important; color: #334155 !important; }
  [data-baseweb="option"]:hover { background: #eef3fb !important; color: #0D4FA3 !important; }

  /* Text area */
  .stTextArea textarea {
    background: #ffffff !important;
    border: 1.5px solid #dde5f0 !important;
    border-radius: 10px !important;
    color: #1e293b !important;
    font-size: 0.85rem !important;
  }
  .stTextArea textarea:focus {
    border-color: #0D4FA3 !important;
    box-shadow: 0 0 0 3px rgba(13,79,163,0.1) !important;
  }
  .stTextArea label { color: #64748b !important; font-size: 0.75rem !important; }

  /* Primary button */
  div[data-testid="stButton"] > button[kind="primary"] {
    background: linear-gradient(135deg, #0D4FA3 0%, #1a6fc4 100%) !important;
    color: white !important;
    border: none !important;
    font-weight: 800 !important;
    font-size: 0.9rem !important;
    letter-spacing: 0.04em !important;
    border-radius: 12px !important;
    padding: 13px 0 !important;
    animation: button-lift 3s ease-in-out infinite !important;
    transition: transform 0.15s ease !important;
  }
  div[data-testid="stButton"] > button[kind="primary"]:hover {
    transform: translateY(-2px) !important;
    background: linear-gradient(135deg, #0a3d8a 0%, #0D4FA3 100%) !important;
  }

  /* Secondary buttons */
  div[data-testid="stButton"] > button[kind="secondary"] {
    background: #ffffff !important;
    border: 1.5px solid #dde5f0 !important;
    color: #475569 !important;
    border-radius: 8px !important;
    font-size: 0.78rem !important;
    font-weight: 600 !important;
    transition: all 0.15s ease !important;
  }
  div[data-testid="stButton"] > button[kind="secondary"]:hover {
    border-color: #0D4FA3 !important;
    color: #0D4FA3 !important;
    background: #eef3fb !important;
  }

  /* Brand profile card */
  .brand-card {
    background: #ffffff;
    border: 1px solid #dde5f0;
    border-left: 4px solid #0D4FA3;
    border-radius: 12px;
    padding: 14px 16px;
    font-size: 0.8rem;
    line-height: 1.65;
    color: #475569;
    box-shadow: 0 2px 8px rgba(13,79,163,0.06);
  }
  .brand-card strong { color: #1e293b; }
  .brand-price-badge {
    display: inline-block;
    background: #eef3fb;
    border: 1px solid #b3c9e8;
    color: #0D4FA3;
    padding: 2px 9px;
    border-radius: 20px;
    font-size: 0.7rem;
    font-weight: 700;
    margin-left: 6px;
  }
  .brand-strength {
    color: #059669;
    font-weight: 600;
    font-size: 0.72rem;
  }

  /* Data badge */
  .data-badge {
    background: #f0fdf4;
    border: 1px solid #86efac;
    border-radius: 10px;
    padding: 10px 14px;
    font-size: 0.82rem;
    color: #15803d;
    margin-bottom: 10px;
    display: flex;
    align-items: flex-start;
    gap: 8px;
  }

  /* Expander */
  [data-testid="stExpander"] {
    background: #ffffff !important;
    border: 1px solid #dde5f0 !important;
    border-radius: 10px !important;
  }
  [data-testid="stExpander"] summary { color: #64748b !important; font-size: 0.8rem !important; }
  [data-testid="stExpander"] summary:hover { color: #0D4FA3 !important; }

  /* Chat messages */
  [data-testid="stChatMessage"] {
    background: #ffffff !important;
    border: 1px solid #dde5f0 !important;
    border-radius: 14px !important;
    margin-bottom: 12px !important;
    box-shadow: 0 2px 8px rgba(13,79,163,0.05) !important;
  }
  [data-testid="stChatMessageContent"] { color: #334155 !important; }
  [data-testid="stChatMessageContent"] h1,
  [data-testid="stChatMessageContent"] h2,
  [data-testid="stChatMessageContent"] h3 { color: #1e293b !important; }
  [data-testid="stChatMessageContent"] strong { color: #0f172a !important; }
  [data-testid="stChatMessageContent"] a { color: #0D4FA3 !important; }

  /* Chat input */
  [data-testid="stChatInput"] {
    background: #ffffff !important;
    border: 1.5px solid #dde5f0 !important;
    border-radius: 14px !important;
  }
  [data-testid="stChatInput"]:focus-within {
    border-color: #0D4FA3 !important;
    box-shadow: 0 0 0 3px rgba(13,79,163,0.1) !important;
  }
  [data-testid="stChatInput"] textarea { color: #1e293b !important; background: transparent !important; }
  [data-testid="stChatInput"] textarea::placeholder { color: #94a3b8 !important; }

  /* Welcome box */
  .welcome-box {
    text-align: center;
    padding: 60px 32px;
    color: #64748b;
  }
  .welcome-icon {
    font-size: 3.5rem;
    margin-bottom: 16px;
  }
  .welcome-box h3 {
    color: #1e293b;
    font-size: 1.3rem;
    margin: 0 0 10px;
    font-weight: 800;
    letter-spacing: -0.02em;
  }
  .welcome-box p {
    font-size: 0.88rem;
    line-height: 1.75;
    margin: 0;
    color: #64748b;
  }
  .welcome-box .start-hint {
    display: inline-block;
    margin-top: 20px;
    background: #eef3fb;
    border: 1px solid #b3c9e8;
    border-radius: 8px;
    padding: 8px 18px;
    font-size: 0.78rem;
    color: #0D4FA3;
    font-weight: 600;
  }

  /* Download button */
  [data-testid="stDownloadButton"] button {
    font-size: 0.75rem !important;
    padding: 5px 14px !important;
    background: #eef3fb !important;
    border: 1px solid #b3c9e8 !important;
    color: #0D4FA3 !important;
    border-radius: 8px !important;
    font-weight: 600 !important;
  }
  [data-testid="stDownloadButton"] button:hover {
    background: #dde8f5 !important;
  }

  /* Dataframe */
  [data-testid="stDataFrame"] { border-radius: 10px !important; overflow: hidden; }

  /* Alerts */
  [data-testid="stAlert"] {
    background: #f0fdf4 !important;
    border: 1px solid #86efac !important;
    border-radius: 10px !important;
    color: #15803d !important;
  }

  /* Divider */
  hr { border-color: #e8edf5 !important; }

  /* File uploader */
  [data-testid="stFileUploadDropzone"] {
    background: #f8fafc !important;
    border: 1.5px dashed #b3c9e8 !important;
    border-radius: 12px !important;
    transition: border-color 0.2s ease !important;
  }
  [data-testid="stFileUploadDropzone"]:hover {
    border-color: #0D4FA3 !important;
    background: #eef3fb !important;
  }
  [data-testid="stFileUploadDropzone"] * { color: #64748b !important; }

  /* Caption */
  .stCaption, [data-testid="stCaptionContainer"] p { color: #94a3b8 !important; font-size: 0.75rem !important; }

  /* Spinner */
  [data-testid="stSpinner"] { color: #0D4FA3 !important; }
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────────
# SESSION STATE
# ─────────────────────────────────────────────────────────────────

_defaults = {
    "messages":      [],
    "_pending":      None,
    "data_summary":  "",
    "data_filename": "",
    "data_df":       None,
    "_show_preview": False,
}
for k, v in _defaults.items():
    if k not in st.session_state:
        st.session_state[k] = v

# ─────────────────────────────────────────────────────────────────
# SIDEBAR — API KEY + CONTROLS
# ─────────────────────────────────────────────────────────────────

with st.sidebar:
    _sb_logo = get_logo_html(height=36)
    if _sb_logo:
        st.markdown(
            f'<div style="padding:6px 0 14px">{_sb_logo}</div>',
            unsafe_allow_html=True,
        )
    else:
        st.markdown(
            '<div style="font-size:0.95rem;font-weight:800;color:#0D4FA3;padding:4px 0 14px;letter-spacing:-0.01em">'
            'Genomma Lab.</div>',
            unsafe_allow_html=True,
        )
    st.markdown("### ⚙️ Settings")
    api_key = st.text_input(
        "Anthropic API Key",
        type="password",
        value=os.environ.get("ANTHROPIC_API_KEY", ""),
        help="Your key stays on your machine and is never sent anywhere else.",
    )
    st.divider()
    if st.button("🗑️ Clear conversation", use_container_width=True):
        st.session_state.messages      = []
        st.session_state._pending      = None
        st.rerun()
    st.caption(f"Today: {datetime.now().strftime('%B %d, %Y')}")
    st.caption("Powered by Claude · Live web search")

# ─────────────────────────────────────────────────────────────────
# CLAUDE CALLER
# ─────────────────────────────────────────────────────────────────

def get_client() -> anthropic.Anthropic:
    key = api_key or os.environ.get("ANTHROPIC_API_KEY", "")
    if not key:
        st.error("⚠️ Enter your Anthropic API key in the sidebar (☰ top-left).")
        st.stop()
    return anthropic.Anthropic(api_key=key)


def call_agent(client: anthropic.Anthropic, api_messages: list) -> str:
    """Run Claude with web search until end_turn."""
    max_rounds = 12
    response   = None
    for _ in range(max_rounds):
        response = client.messages.create(
            model="claude-opus-4-6",
            max_tokens=12000,
            system=SYSTEM_PROMPT,
            tools=[{"type": "web_search_20260209", "name": "web_search"}],
            messages=api_messages,
        )
        if response.stop_reason == "end_turn":
            break
        if response.stop_reason == "pause_turn":
            api_messages.append({"role": "assistant", "content": response.content})
            continue
        break

    texts = [
        b.text for b in (response.content if response else [])
        if getattr(b, "type", None) == "text"
    ]
    return "\n".join(texts) if texts else "No response generated."


# ─────────────────────────────────────────────────────────────────
# HEADER
# ─────────────────────────────────────────────────────────────────

_logo_html = get_logo_html(height=46)
st.markdown(f"""
<div class="hub-header">
  <div style="display:flex;align-items:center;justify-content:space-between;flex-wrap:wrap;gap:12px">
    <div>
      <div class="header-badge">
        <span class="live-dot"></span> LIVE &nbsp;&middot;&nbsp; REAL-TIME INTELLIGENCE
      </div>
      <h1>{"" if _logo_html else ""}Innovation Intelligence Platform</h1>
      <p>Skin Care &nbsp;&middot;&nbsp; Hair Care &nbsp;&middot;&nbsp; Body Care &nbsp;&nbsp;|&nbsp;&nbsp;
         US Market &nbsp;&nbsp;|&nbsp;&nbsp; Claude AI + Live Web Research</p>
    </div>
    {f'<div style="opacity:0.92">{_logo_html}</div>' if _logo_html else ""}
  </div>
  <div class="header-stats">
    <div class="hstat"><span>6</span>BRANDS</div>
    <div class="hstat"><span>3</span>CATEGORIES</div>
    <div class="hstat"><span>6</span>REPORT SECTIONS</div>
  </div>
</div>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────────
# MAIN LAYOUT
# ─────────────────────────────────────────────────────────────────

col_controls, col_results = st.columns([1, 2], gap="large")

# ══════════════════════════════════════════════════════════════════
# LEFT COLUMN — Dashboard Controls
# ══════════════════════════════════════════════════════════════════

with col_controls:

    # ── DROPDOWN SELECTORS ────────────────────────────────────────
    st.markdown('<div class="glass-card">', unsafe_allow_html=True)
    st.markdown('<div class="sec-label">🎛️ Analysis Parameters</div>', unsafe_allow_html=True)

    sel_category = st.selectbox(
        "📂 Category",
        ["Skin Care", "Hair Care", "Body Care"],
        key="sel_category",
    )

    available_brands = CATEGORY_BRANDS.get(sel_category, list(BRAND_PROFILES.keys()))
    sel_brand = st.selectbox(
        "🏷️ Brand",
        available_brands,
        key="sel_brand",
    )

    sel_market = st.selectbox(
        "🎯 Target Market",
        ["US Hispanic", "US General Market", "Both"],
        key="sel_market",
    )

    sel_innovation = st.selectbox(
        "💡 Innovation Type",
        ["New to Market", "Competitive Mimic + Improve", "Both"],
        key="sel_innovation",
    )

    st.markdown('</div>', unsafe_allow_html=True)

    # ── OPTIONAL FOCUS ────────────────────────────────────────────
    with st.expander("✏️ Add focus or context (optional)"):
        extra_context = st.text_area(
            "Additional focus",
            placeholder="e.g. Focus on hair loss for women 40+, or premium positioning for CVS",
            height=80,
            key="extra_context",
            label_visibility="collapsed",
        )

    # ── GENERATE BUTTON ───────────────────────────────────────────
    st.markdown("<br>", unsafe_allow_html=True)
    if st.button("▶ Generate Innovation Report", use_container_width=True, type="primary", key="btn_run"):
        icon  = BRAND_PROFILES.get(sel_brand, {}).get("icon", "")
        label = f"{icon} {sel_brand} — {sel_category} · {sel_market} · {sel_innovation}"
        prompt = build_prompt(sel_category, sel_brand, sel_market, sel_innovation,
                              st.session_state.get("extra_context", ""))
        st.session_state.messages.append({"role": "user", "content": label})
        st.session_state._pending = prompt
        st.rerun()

    st.divider()

    # ── FILE UPLOAD ───────────────────────────────────────────────
    st.markdown('<div class="sec-label">📂 Upload Market Data</div>', unsafe_allow_html=True)
    st.caption("IRI / Nielsen (Excel, CSV) · Presentations (PDF, PPT)")

    if st.session_state.data_filename:
        st.markdown(
            f'<div class="data-badge">'
            f'<span style="font-size:1rem">✅</span>'
            f'<div><strong style="color:#34d399">{st.session_state.data_filename}</strong><br>'
            f'<span style="font-size:0.74rem;color:#059669">Real data active · used in all analyses</span></div>'
            f'</div>',
            unsafe_allow_html=True,
        )
        c1, c2 = st.columns(2)
        with c1:
            if st.button("✕ Remove", use_container_width=True, key="btn_remove"):
                st.session_state.data_summary  = ""
                st.session_state.data_filename = ""
                st.session_state.data_df       = None
                st.session_state._show_preview = False
                st.rerun()
        with c2:
            preview_label = "▲ Hide" if st.session_state._show_preview else "👁 Preview"
            if st.button(preview_label, use_container_width=True, key="btn_preview"):
                st.session_state._show_preview = not st.session_state._show_preview
                st.rerun()
        if st.session_state._show_preview and st.session_state.data_df is not None:
            st.dataframe(st.session_state.data_df.head(25), use_container_width=True)
    else:
        uploaded = st.file_uploader(
            "Drop your file here",
            type=["xlsx", "xls", "csv", "pdf", "pptx", "ppt"],
            key="file_upload",
            label_visibility="collapsed",
        )
        if uploaded:
            with st.spinner(f"Reading {uploaded.name}…"):
                df, summary, ftype = parse_data_file(uploaded)
            if ftype != "error":
                st.session_state.data_df       = df
                st.session_state.data_summary  = summary
                st.session_state.data_filename = uploaded.name
                st.success(f"✅ Loaded **{uploaded.name}**")
                st.rerun()
            else:
                st.error(summary)

    st.divider()

    # ── SELECTED BRAND PROFILE ────────────────────────────────────
    st.markdown('<div class="sec-label">📋 Brand Profile</div>', unsafe_allow_html=True)
    p = BRAND_PROFILES.get(sel_brand, {})
    if p:
        st.markdown(
            f'<div class="brand-card">'
            f'<strong style="font-size:0.95rem">{p.get("icon","")} {sel_brand}</strong>'
            f'<span class="brand-price-badge">{p.get("price_range","")}</span><br><br>'
            f'<span style="color:#64748b;font-size:0.78rem">{p.get("description","")[:220]}…</span><br><br>'
            f'<span class="brand-strength">✓ Strengths: </span>'
            f'<span style="color:#475569;font-size:0.75rem">{p.get("strengths","")}</span>'
            f'</div>',
            unsafe_allow_html=True,
        )


# ══════════════════════════════════════════════════════════════════
# RIGHT COLUMN — Results & Chat
# ══════════════════════════════════════════════════════════════════

with col_results:
    st.markdown('<div class="sec-label">📊 Research Results & Chat</div>', unsafe_allow_html=True)
    st.caption(
        "Configure parameters on the left → Generate Report. "
        "Then refine any section via the chat below."
    )

    # ── Welcome state ─────────────────────────────────────────────
    if not st.session_state.messages and not st.session_state._pending:
        st.markdown("""
        <div class="welcome-box">
            <div class="welcome-icon">🔬</div>
            <h3>Ready to discover your next launch</h3>
            <p>
                Select a <strong style="color:#1e293b">Category</strong>,
                <strong style="color:#1e293b">Brand</strong>,
                <strong style="color:#1e293b">Target Market</strong>, and
                <strong style="color:#1e293b">Innovation Type</strong> on the left.<br><br>
                Hit <strong style="color:#0D4FA3">Generate Innovation Report</strong> for a full<br>
                6-section intelligence brief with live market data<br>
                and AI-rendered product concept images.
            </p>
            <div class="start-hint">← Configure parameters &amp; click Generate</div>
        </div>
        """, unsafe_allow_html=True)

    # ── Display conversation ──────────────────────────────────────
    for i, msg in enumerate(st.session_state.messages):
        avatar = "🔬" if msg["role"] == "assistant" else "👤"
        with st.chat_message(msg["role"], avatar=avatar):
            if msg["role"] == "assistant":
                st.markdown(msg["content"])
                st.download_button(
                    "📥 Download report",
                    data=msg["content"],
                    file_name=f"innovation_report_{datetime.now().strftime('%Y%m%d_%H%M')}_{i}.txt",
                    mime="text/plain",
                    key=f"dl_{i}",
                )
            else:
                st.markdown(msg["content"])

    # ── Process pending report ────────────────────────────────────
    if st.session_state._pending:
        pending                  = st.session_state._pending
        st.session_state._pending = None

        client = get_client()
        with st.chat_message("assistant", avatar="🔬"):
            with st.spinner(
                "🔍 Searching the web and building your report…  \n"
                "Full analyses typically take 45–90 seconds."
            ):
                api_msgs = []
                for m in st.session_state.messages[:-1]:
                    api_msgs.append({"role": m["role"], "content": m["content"]})
                api_msgs.append({"role": "user", "content": pending})
                result = call_agent(client, api_msgs)

            st.markdown(result)
            st.download_button(
                "📥 Download report",
                data=result,
                file_name=f"innovation_report_{datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                mime="text/plain",
                key="dl_pending",
            )

        st.session_state.messages.append({"role": "assistant", "content": result})
        st.rerun()

    # ── Chat input ────────────────────────────────────────────────
    user_input = st.chat_input(
        "Refine or go deeper… e.g. 'Adjust pricing for CVS' · 'Go deeper on concept 2' · 'Add Hispanic insights'"
    )

    if user_input:
        st.session_state.messages.append({"role": "user", "content": user_input})
        with st.chat_message("user", avatar="👤"):
            st.markdown(user_input)

        client = get_client()
        with st.chat_message("assistant", avatar="🔬"):
            with st.spinner("Researching…"):
                api_msgs = [
                    {"role": m["role"], "content": m["content"]}
                    for m in st.session_state.messages
                ]
                result = call_agent(client, api_msgs)
            st.markdown(result)
            st.download_button(
                "📥 Download report",
                data=result,
                file_name=f"innovation_{datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                mime="text/plain",
                key=f"dl_chat_{len(st.session_state.messages)}",
            )

        st.session_state.messages.append({"role": "assistant", "content": result})
